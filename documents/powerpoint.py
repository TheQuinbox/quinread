from .base import BaseDocument
from pptx import Presentation

class PptxDocument(BaseDocument):
	def __init__(self, path):
		self.path = path

	def read(self):
		self.document = Presentation(self.path)
		final = ""
		for slide in self.document.slides:
			for shape in slide.shapes:
				if hasattr(shape, "text"):
					final += f"{shape.text}\n"
			final += "\n"
		return final

	def close(self):
		super().close()
